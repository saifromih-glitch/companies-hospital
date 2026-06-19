"""
Tool Guard — post-processing layer for tool execution
Catches JSON tool calls in Romih's output and executes them
Does NOT modify core agent logic — sits between agent and user
"""
import json, re, io, os

class ToolGuard:
    """مسك JSON tool calls وتنفيذها"""
    
    def __init__(self):
        self.tools = {
            "create_xlsx": self._create_xlsx,
            "create_csv": self._create_csv,
            "create_docx": self._create_docx,
            "create_pptx": self._create_pptx,
            "create_pdf": self._create_pdf,
            "generate_chart": self._generate_chart,
        }
    
    def scan(self, response: str) -> dict:
        """
        Scan response for tool calls
        Returns: {has_tool: bool, cleaned_text: str, files: [{name, content, mime}]}
        """
        files = []
        cleaned = response
        
        # Pattern 1: JSON tool call block
        json_match = re.search(r'```json\s*\n(.*?)```', response, re.DOTALL)
        if json_match:
            try:
                data = json.loads(json_match.group(1))
                if isinstance(data, dict) and "tool" in data:
                    result = self._execute(data["tool"], data.get("args", {}))
                    if result:
                        files.append(result)
                        cleaned = re.sub(r'```json\s*\n.*?```', '', response, flags=re.DOTALL).strip()
            except:
                pass
        
        # Pattern 2: Inline FILE: marker
        file_match = re.search(r'FILE:(\w+):(\{.*?\})', response, re.DOTALL)
        if file_match:
            try:
                tool = file_match.group(1)
                args = json.loads(file_match.group(2))
                result = self._execute(tool, args)
                if result:
                    files.append(result)
                    cleaned = re.sub(r'FILE:\w+:\{.*?\}', '', response, flags=re.DOTALL).strip()
            except:
                pass
        
        has_tool = len(files) > 0
        return {"has_tool": has_tool, "cleaned_text": cleaned, "files": files}
    
    def _execute(self, tool: str, args: dict) -> dict:
        """Execute a tool and return file info"""
        handler = self.tools.get(tool)
        if not handler:
            return None
        try:
            return handler(args)
        except Exception as e:
            print(f"ToolGuard: {tool} failed: {e}")
            return None
    
    def _create_xlsx(self, args: dict) -> dict:
        """Create actual XLSX file"""
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = args.get("sheet_name", "Sheet1")
        
        title = args.get("title", "")
        if title:
            ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=8)
            cell = ws.cell(row=1, column=1, value=title)
            cell.font = Font(bold=True, size=16)
            cell.alignment = Alignment(horizontal='center')
        
        headers = args.get("headers", [])
        header_row = 2 if title else 1
        header_fill = PatternFill(start_color="3498db", end_color="3498db", fill_type="solid")
        
        for i, h in enumerate(headers, 1):
            cell = ws.cell(row=header_row, column=i, value=h)
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal='center')
        
        data = args.get("data", [])
        for r, row in enumerate(data, header_row + 1):
            for c, val in enumerate(row, 1):
                ws.cell(row=r, column=c, value=val)
        
        # Auto-width
        for col in ws.columns:
            max_length = 0
            for cell in col:
                if cell.value:
                    max_length = max(max_length, len(str(cell.value)))
            ws.column_dimensions[col[0].column_letter].width = min(max_length + 4, 30)
        
        output = io.BytesIO()
        wb.save(output)
        output.seek(0)
        
        return {
            "name": args.get("filename", "romih_file.xlsx"),
            "content": output.getvalue(),
            "mime": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        }
    
    def _create_csv(self, args: dict) -> dict:
        """Create CSV file from data"""
        import csv
        output = io.StringIO()
        writer = csv.writer(output)
        
        if args.get("headers"):
            writer.writerow(args["headers"])
        for row in args.get("data", []):
            writer.writerow(row)
        
        content = output.getvalue()
        return {
            "name": args.get("filename", "romih_data.csv"),
            "content": content.encode('utf-8'),
            "mime": "text/csv"
        }
    
    def _create_docx(self, args: dict) -> dict:
        """Create DOCX file"""
        from docx import Document
        doc = Document()
        
        if args.get("title"):
            doc.add_heading(args["title"], level=1)
        
        for section in args.get("sections", []):
            if section.get("heading"):
                doc.add_heading(section["heading"], level=2)
            if section.get("text"):
                doc.add_paragraph(section["text"])
            if section.get("items"):
                for item in section["items"]:
                    doc.add_paragraph(item, style='List Bullet')
        
        output = io.BytesIO()
        doc.save(output)
        output.seek(0)
        
        return {
            "name": args.get("filename", "romih_doc.docx"),
            "content": output.getvalue(),
            "mime": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        }
    
    def _create_pptx(self, args: dict) -> dict:
        """Create PPTX presentation"""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        title = args.get("title", "Presentation")
        slides = args.get("slides", [])
        
        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        if slides:
            slide.placeholders[1].text = slides[0].get("subtitle", "")
        
        # Content slides
        for s in slides[1:]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = s.get("title", "")
            body = slide.placeholders[1]
            body.text = s.get("content", "")
            if s.get("items"):
                for item in s["items"]:
                    body.text += "\n" + "• " + item
        
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        return {
            "name": args.get("filename", "romih_presentation.pptx"),
            "content": output.getvalue(),
            "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        }
    
    def _create_pdf(self, args: dict) -> dict:
        """Create PDF using fpdf"""
        from fpdf import FPDF
        
        pdf = FPDF()
        pdf.add_page()
        
        # Arabic font support
        try:
            pdf.add_font('Arabic', '', r'C:\Windows\Fonts\arial.ttf', uni=True)
            pdf.set_font('Arabic', '', 12)
        except:
            pdf.set_font('Arial', '', 12)
        
        if args.get("title"):
            pdf.set_font_size(18)
            pdf.cell(0, 10, args["title"], ln=True, align='C')
            pdf.ln(10)
        
        pdf.set_font_size(12)
        for section in args.get("sections", []):
            if section.get("heading"):
                pdf.set_font_size(14)
                pdf.cell(0, 8, section["heading"], ln=True)
                pdf.set_font_size(12)
            if section.get("text"):
                pdf.multi_cell(0, 6, section["text"])
            if section.get("items"):
                for item in section["items"]:
                    pdf.cell(10)
                    pdf.cell(0, 6, "• " + item, ln=True)
            pdf.ln(4)
        
        output = io.BytesIO()
        pdf.output(output)
        output.seek(0)
        
        return {
            "name": args.get("filename", "romih_document.pdf"),
            "content": output.getvalue(),
            "mime": "application/pdf"
        }
    
    def _generate_chart(self, args: dict) -> dict:
        """Generate chart as HTML for Telegram"""
        chart_type = args.get("type", "bar")
        labels = json.dumps(args.get("labels", []))
        values = json.dumps(args.get("values", []))
        title = args.get("title", "Chart")
        
        html = f"""<!DOCTYPE html><html><head><meta charset="UTF-8">
<script src="https://cdn.jsdelivr.net/npm/chart.js"></script>
</head><body>
<canvas id="c"></canvas>
<script>
new Chart(document.getElementById('c'), {{
    type: '{chart_type}',
    data: {{ labels: {labels}, datasets: [{{ label: '{title}', data: {values} }}] }},
    options: {{ responsive: true }}
}});
</script></body></html>"""
        
        return {
            "name": f"{title}.html",
            "content": html.encode('utf-8'),
            "mime": "text/html"
        }

# Global instance
guard = ToolGuard()
